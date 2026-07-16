from fastapi import FastAPI, Request, Form, UploadFile, File, WebSocket, WebSocketDisconnect, Depends, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from fastapi.responses import HTMLResponse, JSONResponse, RedirectResponse, StreamingResponse
from jose import JWTError, jwt
from passlib.context import CryptContext
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from sklearn.ensemble import IsolationForest
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
import json
import sqlite3
import hashlib
from datetime import datetime, timedelta
from io import BytesIO, StringIO
import os
import asyncio
import openai

# ==================== CONFIG ====================
SECRET_KEY = os.environ.get("SECRET_KEY", "metricsai-dev-key-change-in-production")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24 * 7
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")

app = FastAPI(title="MetricsAI v4.1")
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

# ==================== DB ====================
def get_db():
    conn = sqlite3.connect("metricsai.db")
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db()
    c = conn.cursor()
    c.execute("CREATE TABLE IF NOT EXISTS shared_dashboards (id TEXT PRIMARY KEY, created_at TEXT, data_json TEXT, white_label TEXT)")
    c.execute("CREATE TABLE IF NOT EXISTS uploads (id INTEGER PRIMARY KEY AUTOINCREMENT, filename TEXT, source_type TEXT, upload_date TEXT, rows INTEGER, cols INTEGER, user_id INTEGER)")
    c.execute("CREATE TABLE IF NOT EXISTS custom_reports (id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, config_json TEXT, created_at TEXT, user_id INTEGER)")
    c.execute("CREATE TABLE IF NOT EXISTS api_keys (id INTEGER PRIMARY KEY AUTOINCREMENT, service TEXT, key_value TEXT, created_at TEXT)")
    c.execute("CREATE TABLE IF NOT EXISTS users (id INTEGER PRIMARY KEY AUTOINCREMENT, email TEXT UNIQUE NOT NULL, username TEXT UNIQUE NOT NULL, hashed_password TEXT NOT NULL, full_name TEXT, is_active INTEGER DEFAULT 1, created_at TEXT DEFAULT CURRENT_TIMESTAMP)")
    c.execute("CREATE TABLE IF NOT EXISTS user_dashboards (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER NOT NULL, name TEXT NOT NULL, data_json TEXT, white_label TEXT, created_at TEXT DEFAULT CURRENT_TIMESTAMP, FOREIGN KEY (user_id) REFERENCES users (id))")
    c.execute("CREATE TABLE IF NOT EXISTS exports (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER, export_type TEXT, filename TEXT, created_at TEXT DEFAULT CURRENT_TIMESTAMP)")
    c.execute("CREATE TABLE IF NOT EXISTS alerts (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER, metric TEXT, threshold_type TEXT, threshold_value REAL, webhook_url TEXT, is_active INTEGER DEFAULT 1, created_at TEXT DEFAULT CURRENT_TIMESTAMP)")
    c.execute("CREATE TABLE IF NOT EXISTS ai_conversations (id INTEGER PRIMARY KEY AUTOINCREMENT, session_id TEXT, user_message TEXT, ai_response TEXT, created_at TEXT DEFAULT CURRENT_TIMESTAMP)")
    conn.commit()
    conn.close()

init_db()

# ==================== AUTH ====================
def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    return pwd_context.hash(password)

def create_access_token(data: dict, expires_delta: timedelta = None):
    to_encode = data.copy()
    expire = datetime.utcnow() + (expires_delta or timedelta(minutes=15))
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def get_current_user(request: Request):
    token = request.cookies.get("access_token")
    if not token:
        return None
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username = payload.get("sub")
        if not username:
            return None
        conn = get_db()
        c = conn.cursor()
        c.execute("SELECT * FROM users WHERE username = ?", (username,))
        user = c.fetchone()
        conn.close()
        return dict(user) if user else None
    except JWTError:
        return None

# ==================== WEBSOCKET ====================
class ConnectionManager:
    def __init__(self):
        self.active_connections = []
    async def connect(self, websocket: WebSocket):
        await websocket.accept()
        self.active_connections.append(websocket)
    def disconnect(self, websocket: WebSocket):
        if websocket in self.active_connections:
            self.active_connections.remove(websocket)
    async def broadcast(self, message: dict):
        disconnected = []
        for conn in self.active_connections:
            try:
                await conn.send_json(message)
            except:
                disconnected.append(conn)
        for conn in disconnected:
            self.disconnect(conn)

manager = ConnectionManager()

async def broadcast_metrics_update(session_id: str):
    data = get_session_data(session_id)
    if data:
        await manager.broadcast({"type": "metrics_update", "session_id": session_id, "metrics": data["metrics"], "timestamp": datetime.now().isoformat()})

# ==================== THEME ====================
PLOTLY_THEME = {
    "template": "plotly_dark",
    "paper_bgcolor": "rgba(0,0,0,0)",
    "plot_bgcolor": "rgba(0,0,0,0)",
    "font_color": "#94A3B8",
    "font_family": "Inter, sans-serif",
    "margin": dict(l=10, r=10, t=30, b=10),
    "showlegend": False,
    "xaxis": dict(showgrid=True, gridcolor="rgba(255,255,255,0.03)", zeroline=False),
    "yaxis": dict(showgrid=True, gridcolor="rgba(255,255,255,0.03)", zeroline=False),
}

COLORS = {
    "primary": "#F8FAFC", "secondary": "#94A3B8", "accent": "#34D399",
    "blue": "#60A5FA", "purple": "#BF5AF2", "orange": "#FF9F0A",
    "red": "#FF453A", "teal": "#5AC8FA", "indigo": "#5E5CE6",
    "palette": ["#34D399", "#60A5FA", "#BF5AF2", "#FF9F0A", "#5AC8FA", "#5E5CE6", "#FF453A", "#F8FAFC"]
}

def apply_theme(fig):
    fig.update_layout(**PLOTLY_THEME)
    return fig

# ==================== DATA ====================
def generate_sample_data():
    np.random.seed(42)
    dates = pd.date_range("2023-01-01", "2025-06-30", freq="D")
    n = len(dates)
    users = [f"user_{i:04d}" for i in range(500)]
    channels = ["organic", "instagram", "google", "facebook", "email", "tiktok", "referral"]
    events = ["impression", "click", "visit", "add_to_cart", "purchase"]
    df = pd.DataFrame({
        "date": np.random.choice(dates, n),
        "user_id": np.random.choice(users, n),
        "revenue": np.where(np.random.random(n) > 0.7, np.random.exponential(200, n).round(2), 0),
        "channel": np.random.choice(channels, n, p=[0.25, 0.15, 0.20, 0.15, 0.10, 0.10, 0.05]),
        "event_type": np.random.choice(events, n, p=[0.35, 0.20, 0.20, 0.10, 0.15]),
        "marketing_cost": np.random.choice([0, 0, 0, 30, 50, 80, 120, 200], n),
        "subscription": np.random.choice(["monthly", "annual", "none"], n, p=[0.15, 0.05, 0.80]),
        "session_duration": np.random.exponential(300, n).round(0),
        "page_views": np.random.poisson(5, n),
        "device": np.random.choice(["desktop", "mobile", "tablet"], n, p=[0.5, 0.4, 0.1]),
        "country": np.random.choice(["RU", "US", "DE", "FR", "UK"], n, p=[0.4, 0.25, 0.15, 0.10, 0.10]),
    })
    df["date"] = pd.to_datetime(df["date"])
    return df.sort_values("date").reset_index(drop=True)

# ==================== METRICS ====================
def calculate_metrics(df):
    metrics = {}
    metrics["total_users"] = int(df["user_id"].nunique()) if "user_id" in df.columns else len(df)
    metrics["total_revenue"] = float(df["revenue"].sum()) if "revenue" in df.columns else 0.0
    metrics["total_orders"] = int(len(df[df["revenue"] > 0])) if "revenue" in df.columns else len(df)
    metrics["arpu"] = metrics["total_revenue"] / metrics["total_users"] if metrics["total_users"] > 0 else 0.0
    metrics["aov"] = metrics["total_revenue"] / metrics["total_orders"] if metrics["total_orders"] > 0 else 0.0

    if "channel" in df.columns and "marketing_cost" in df.columns:
        channel_stats = df.groupby("channel").agg({"user_id": "nunique", "marketing_cost": "sum", "revenue": "sum"}).reset_index()
        channel_stats["cac"] = channel_stats["marketing_cost"] / channel_stats["user_id"]
        channel_stats["roi"] = (channel_stats["revenue"] - channel_stats["marketing_cost"]) / channel_stats["marketing_cost"]
        metrics["cac_by_channel"] = channel_stats.to_dict("records")
        metrics["avg_cac"] = float(channel_stats["cac"].mean())
        metrics["best_channel"] = channel_stats.loc[channel_stats["roi"].idxmax(), "channel"] if not channel_stats.empty else "N/A"
        metrics["worst_channel"] = channel_stats.loc[channel_stats["roi"].idxmin(), "channel"] if not channel_stats.empty else "N/A"
    else:
        metrics["avg_cac"] = 0.0
        metrics["cac_by_channel"] = []
        metrics["best_channel"] = "N/A"
        metrics["worst_channel"] = "N/A"

    if "date" in df.columns and "user_id" in df.columns and "revenue" in df.columns:
        user_revenue = df.groupby("user_id")["revenue"].sum()
        metrics["ltv"] = float(user_revenue.mean() * 2.0)
        metrics["ltv_p90"] = float(user_revenue.quantile(0.90))
        metrics["ltv_p50"] = float(user_revenue.quantile(0.50))
    else:
        metrics["ltv"] = metrics["arpu"] * 3.0
        metrics["ltv_p90"] = metrics["ltv"]
        metrics["ltv_p50"] = metrics["ltv"]

    metrics["ltv_cac_ratio"] = metrics["ltv"] / metrics["avg_cac"] if metrics["avg_cac"] > 0 else 0.0

    if "session_duration" in df.columns:
        metrics["avg_session_duration"] = float(df["session_duration"].mean())
        metrics["avg_page_views"] = float(df["page_views"].mean())
    else:
        metrics["avg_session_duration"] = 0
        metrics["avg_page_views"] = 0

    if "device" in df.columns:
        metrics["device_breakdown"] = df["device"].value_counts().to_dict()
    if "country" in df.columns:
        metrics["country_breakdown"] = df.groupby("country")["revenue"].sum().to_dict()

    if "date" in df.columns:
        df["month_str"] = df["date"].dt.strftime("%Y-%m")
        monthly = df.groupby("month_str").agg({"revenue": "sum", "user_id": "nunique", "marketing_cost": "sum"}).reset_index()
        monthly["arpu"] = monthly["revenue"] / monthly["user_id"]
        monthly["cac"] = monthly["marketing_cost"] / monthly["user_id"]
        monthly["profit"] = monthly["revenue"] - monthly["marketing_cost"]
        monthly = monthly.rename(columns={"month_str": "month"})
        metrics["monthly_trend"] = monthly.to_dict("records")
        if len(monthly) >= 2:
            metrics["revenue_growth"] = float((monthly["revenue"].iloc[-1] - monthly["revenue"].iloc[-2]) / monthly["revenue"].iloc[-2] * 100)
            metrics["user_growth"] = float((monthly["user_id"].iloc[-1] - monthly["user_id"].iloc[-2]) / monthly["user_id"].iloc[-2] * 100)
        else:
            metrics["revenue_growth"] = 0
            metrics["user_growth"] = 0
    else:
        metrics["monthly_trend"] = []
        metrics["revenue_growth"] = 0
        metrics["user_growth"] = 0

    if "event_type" in df.columns:
        funnel_order = ["impression", "click", "visit", "add_to_cart", "purchase"]
        funnel_counts = df["event_type"].value_counts().to_dict()
        funnel_data = []
        prev_count = None
        for step in funnel_order:
            if step in funnel_counts:
                count = int(funnel_counts[step])
                conversion = (count / prev_count * 100) if prev_count and prev_count > 0 else 100
                funnel_data.append({"step": step.upper(), "users": count, "conversion": round(conversion, 1)})
                prev_count = count
        if not funnel_data:
            for k, v in funnel_counts.items():
                funnel_data.append({"step": k.upper(), "users": int(v), "conversion": 100})
        metrics["funnel"] = funnel_data
        metrics["overall_conversion"] = round(funnel_data[-1]["users"] / funnel_data[0]["users"] * 100, 2) if funnel_data else 0
    else:
        metrics["funnel"] = []
        metrics["overall_conversion"] = 0

    if "date" in df.columns and "user_id" in df.columns:
        df_sorted = df.sort_values("date")
        first_seen = df_sorted.groupby("user_id")["date"].min().reset_index()
        first_seen.columns = ["user_id", "first_date"]
        df_merged = df.merge(first_seen, on="user_id")
        df_merged["days_since_first"] = (df_merged["date"] - df_merged["first_date"]).dt.days
        cohorts = []
        for day in [0, 1, 7, 14, 30, 60, 90]:
            active_users = df_merged[df_merged["days_since_first"] == day]["user_id"].nunique()
            total_users = first_seen["user_id"].nunique()
            retention = (active_users / total_users * 100) if total_users > 0 else 0
            cohorts.append({"day": day, "retention": round(retention, 1)})
        metrics["retention"] = cohorts
        metrics["d30_retention"] = round(cohorts[4]["retention"] if len(cohorts) > 4 else 0, 1)
    else:
        metrics["retention"] = []
        metrics["d30_retention"] = 0

    return metrics

# ==================== ANOMALY DETECTION ====================
def detect_anomalies(df):
    if df.empty or "date" not in df.columns or "revenue" not in df.columns:
        return []
    daily = df.groupby(df["date"].dt.date)["revenue"].sum().reset_index()
    daily["date"] = pd.to_datetime(daily["date"])
    daily = daily.sort_values("date").reset_index(drop=True)
    if len(daily) < 10:
        return []
    daily["dow"] = daily["date"].dt.dayofweek
    daily["month"] = daily["date"].dt.month
    daily["revenue_lag1"] = daily["revenue"].shift(1).fillna(daily["revenue"].mean())
    daily["revenue_rolling7"] = daily["revenue"].rolling(7, min_periods=1).mean()
    features = daily[["revenue", "dow", "month", "revenue_lag1", "revenue_rolling7"]].values
    scaler = StandardScaler()
    features_scaled = scaler.fit_transform(features)
    clf = IsolationForest(contamination=0.05, random_state=42)
    daily["anomaly"] = clf.fit_predict(features_scaled)
    daily["anomaly_score"] = clf.decision_function(features_scaled)
    anomalies = daily[daily["anomaly"] == -1].to_dict("records")
    return anomalies

# ==================== RFM (FIXED) ====================
def calculate_rfm(df):
    if "user_id" not in df.columns or "date" not in df.columns or "revenue" not in df.columns:
        return None
    now = df["date"].max()
    rfm = df.groupby("user_id", as_index=False).agg(
        recency=("date", lambda x: (now - x.max()).days),
        frequency=("revenue", "count"),
        monetary=("revenue", "sum")
    )
    features = rfm[["recency", "frequency", "monetary"]].values
    scaler = StandardScaler()
    features_scaled = scaler.fit_transform(features)
    kmeans = KMeans(n_clusters=4, random_state=42, n_init=10)
    rfm["segment"] = kmeans.fit_predict(features_scaled)
    seg_monetary = rfm.groupby("segment")["monetary"].mean().sort_values(ascending=False)
    name_map = {
        seg_monetary.index[0]: "Champions",
        seg_monetary.index[1]: "Loyal",
        seg_monetary.index[2]: "At Risk",
        seg_monetary.index[3]: "Hibernating"
    }
    rfm["segment_name"] = rfm["segment"].map(name_map)
    return rfm.to_dict("records")

# ==================== CHARTS ====================
def build_charts(df, metrics):
    charts = {}
    if metrics.get("monthly_trend"):
        monthly_df = pd.DataFrame(metrics["monthly_trend"])
        fig = px.line(monthly_df, x="month", y="revenue", template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(line_color=COLORS["primary"], line_width=2.5, fill="tozeroy", fillcolor="rgba(248,250,252,0.05)")
        fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25))
        charts["revenue"] = fig.to_html(full_html=False, include_plotlyjs="cdn")

        fig2 = px.area(monthly_df, x="month", y="profit", template="plotly_dark", color_discrete_sequence=[COLORS["accent"]])
        fig2 = apply_theme(fig2)
        fig2.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25))
        charts["profit"] = fig2.to_html(full_html=False, include_plotlyjs=False)

        fig3 = px.bar(monthly_df, x="month", y="arpu", template="plotly_dark")
        fig3 = apply_theme(fig3)
        fig3.update_traces(marker_color=COLORS["secondary"], marker_opacity=0.6, marker_line_width=0)
        fig3.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25))
        charts["arpu"] = fig3.to_html(full_html=False, include_plotlyjs=False)

        fig4 = make_subplots(specs=[[{"secondary_y": True}]])
        fig4.add_trace(go.Scatter(x=monthly_df["month"], y=monthly_df["cac"], name="CAC", line=dict(color=COLORS["red"], width=2)), secondary_y=False)
        fig4.add_trace(go.Scatter(x=monthly_df["month"], y=monthly_df["arpu"]*3, name="LTV (est.)", line=dict(color=COLORS["accent"], width=2)), secondary_y=True)
        fig4 = apply_theme(fig4)
        fig4.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25), showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, font=dict(size=10)))
        charts["cac_ltv"] = fig4.to_html(full_html=False, include_plotlyjs=False)

    if metrics.get("funnel"):
        funnel_df = pd.DataFrame(metrics["funnel"])
        fig = px.funnel(funnel_df, x="users", y="step", template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(marker_color=COLORS["blue"], textposition="inside", textinfo="value+percent initial", textfont=dict(size=11, color="white"))
        fig.update_layout(height=280, margin=dict(l=10, r=10, t=10, b=10))
        charts["funnel"] = fig.to_html(full_html=False, include_plotlyjs=False)

    if "channel" in df.columns:
        channel_df = df.groupby("channel")["revenue"].sum().reset_index().sort_values("revenue", ascending=True)
        fig = px.bar(channel_df, x="revenue", y="channel", orientation="h", template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(marker_color=COLORS["primary"], marker_opacity=0.8, texttemplate="%{x:,.0f}", textposition="outside", textfont=dict(size=10))
        fig.update_layout(height=220, margin=dict(l=10, r=40, t=10, b=10))
        charts["channel"] = fig.to_html(full_html=False, include_plotlyjs=False)

        if metrics.get("cac_by_channel"):
            cac_df = pd.DataFrame(metrics["cac_by_channel"]).sort_values("cac", ascending=True)
            fig = px.bar(cac_df, x="cac", y="channel", orientation="h", template="plotly_dark")
            fig = apply_theme(fig)
            fig.update_traces(marker_color=COLORS["orange"], marker_opacity=0.8, texttemplate="%{x:.0f}", textposition="outside")
            fig.update_layout(height=220, margin=dict(l=10, r=40, t=10, b=10))
            charts["cac_channel"] = fig.to_html(full_html=False, include_plotlyjs=False)

    if metrics.get("device_breakdown"):
        device_df = pd.DataFrame(list(metrics["device_breakdown"].items()), columns=["device", "count"])
        fig = px.pie(device_df, values="count", names="device", template="plotly_dark", hole=0.6)
        fig = apply_theme(fig)
        fig.update_traces(marker_colors=COLORS["palette"][:3], textinfo="percent", textfont=dict(size=11), hovertemplate="%{label}: %{value:,}<extra></extra>")
        fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=10), showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=-0.15, xanchor="center", x=0.5, font=dict(size=10)))
        charts["device"] = fig.to_html(full_html=False, include_plotlyjs=False)

    if metrics.get("country_breakdown"):
        country_df = pd.DataFrame(list(metrics["country_breakdown"].items()), columns=["country", "revenue"]).sort_values("revenue", ascending=False)
        fig = px.bar(country_df, x="country", y="revenue", template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(marker_color=COLORS["purple"], marker_opacity=0.8, texttemplate="%{y:,.0f}", textposition="outside")
        fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25))
        charts["country"] = fig.to_html(full_html=False, include_plotlyjs=False)

    if metrics.get("retention"):
        ret_df = pd.DataFrame(metrics["retention"])
        fig = px.line(ret_df, x="day", y="retention", template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(line_color=COLORS["teal"], line_width=2.5, fill="tozeroy", fillcolor="rgba(90,200,250,0.1)")
        fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25), xaxis_title="Days", yaxis_title="Retention %")
        charts["retention"] = fig.to_html(full_html=False, include_plotlyjs=False)

    if "date" in df.columns:
        df["dow"] = df["date"].dt.day_name()
        df["hour"] = df["date"].dt.hour
        heatmap_df = df.groupby(["dow", "hour"])["revenue"].sum().reset_index()
        if not heatmap_df.empty:
            pivot = heatmap_df.pivot(index="dow", columns="hour", values="revenue").fillna(0)
            dow_order = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]
            pivot = pivot.reindex([d for d in dow_order if d in pivot.index])
            fig = px.imshow(pivot, template="plotly_dark", color_continuous_scale=["rgba(0,0,0,0)", COLORS["accent"]])
            fig = apply_theme(fig)
            fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25), coloraxis_showscale=False)
            charts["heatmap"] = fig.to_html(full_html=False, include_plotlyjs=False)

    if "session_duration" in df.columns:
        fig = px.histogram(df, x="session_duration", nbins=30, template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(marker_color=COLORS["indigo"], marker_opacity=0.7)
        fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25), xaxis_title="Seconds", showlegend=False)
        charts["session_dist"] = fig.to_html(full_html=False, include_plotlyjs=False)

    anomalies = detect_anomalies(df)
    if anomalies and metrics.get("monthly_trend"):
        monthly_df = pd.DataFrame(metrics["monthly_trend"])
        fig = px.line(monthly_df, x="month", y="revenue", template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(line_color=COLORS["primary"], line_width=2.5)
        anomaly_dates = [str(a["date"])[:7] for a in anomalies]
        anomaly_vals = [a["revenue"] for a in anomalies]
        fig.add_trace(go.Scatter(x=anomaly_dates, y=anomaly_vals, mode="markers", name="Anomaly", marker=dict(color=COLORS["red"], size=12, symbol="diamond")))
        fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25), showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, font=dict(size=10)))
        charts["anomalies"] = fig.to_html(full_html=False, include_plotlyjs=False)

    return charts

def calculate_revenue_metrics(df):
    rev = {}
    if df.empty or "date" not in df.columns or "revenue" not in df.columns:
        rev["mrr"] = 0; rev["arr"] = 0; rev["mrr_growth"] = 0
        rev["waterfall"] = []; rev["ltv_curve"] = []; rev["subscription"] = []
        rev["cohort"] = []; rev["forecast"] = []
        return rev
    df["month_str"] = df["date"].dt.strftime("%Y-%m").astype(str)
    monthly_rev = df.groupby("month_str")["revenue"].sum().reset_index()
    monthly_rev["month"] = monthly_rev["month_str"].astype(str)
    monthly_rev = monthly_rev[["month", "revenue"]]
    rev["mrr"] = float(monthly_rev["revenue"].iloc[-1]) if len(monthly_rev) > 0 else 0
    rev["arr"] = rev["mrr"] * 12
    rev["mrr_growth"] = float((monthly_rev["revenue"].iloc[-1] - monthly_rev["revenue"].iloc[-2]) / monthly_rev["revenue"].iloc[-2] * 100) if len(monthly_rev) >= 2 and monthly_rev["revenue"].iloc[-2] > 0 else 0
    rev["waterfall"] = monthly_rev.to_dict("records")
    if "user_id" in df.columns:
        user_ltv = df.groupby("user_id")["revenue"].sum().sort_values(ascending=False)
        rev["ltv_curve"] = [{"percentile": i, "ltv": float(user_ltv.quantile(i/100))} for i in range(0, 101, 5)]
    else:
        rev["ltv_curve"] = []
    if "subscription" in df.columns:
        sub_breakdown = df.groupby("subscription")["revenue"].sum().reset_index().to_dict("records")
        rev["subscription"] = sub_breakdown
    else:
        rev["subscription"] = []
    if "user_id" in df.columns and "date" in df.columns:
        first_purchase = df.groupby("user_id")["date"].min().reset_index(name="first_date")
        first_purchase["cohort"] = first_purchase["first_date"].dt.to_period("M").astype(str)
        df_c = df.merge(first_purchase[["user_id", "cohort"]], on="user_id")
        df_c["periods"] = (df_c["date"].dt.to_period("M") - df_c["cohort"].apply(lambda x: pd.Period(x, freq="M"))).apply(lambda x: x.n)
        cohort_data = df_c.groupby(["cohort", "periods"])["user_id"].nunique().reset_index()
        cohort_sizes = first_purchase.groupby("cohort")["user_id"].nunique().reset_index(name="size")
        cohort_data = cohort_data.merge(cohort_sizes, on="cohort")
        cohort_data["retention"] = (cohort_data["user_id"] / cohort_data["size"] * 100).round(1)
        rev["cohort"] = cohort_data[cohort_data["periods"] <= 6].to_dict("records")
    else:
        rev["cohort"] = []
    if len(monthly_rev) >= 3:
        last3 = monthly_rev["revenue"].iloc[-3:].mean()
        rev["forecast"] = [
            {"month": "Forecast 1M", "revenue": round(last3 * 1.05, 2), "type": "forecast"},
            {"month": "Forecast 2M", "revenue": round(last3 * 1.10, 2), "type": "forecast"},
            {"month": "Forecast 3M", "revenue": round(last3 * 1.15, 2), "type": "forecast"},
        ]
    else:
        rev["forecast"] = []
    return rev

def build_revenue_charts(df, rev_metrics):
    charts = {}
    if rev_metrics.get("waterfall") and len(rev_metrics["waterfall"]) > 0:
        try:
            water_df = pd.DataFrame(rev_metrics["waterfall"])
            if len(water_df) > 0 and "month" in water_df.columns and "revenue" in water_df.columns:
                water_df["month"] = water_df["month"].astype(str)
                water_df["revenue"] = pd.to_numeric(water_df["revenue"], errors="coerce").fillna(0)
                fig = px.area(water_df, x="month", y="revenue", template="plotly_dark", color_discrete_sequence=[COLORS["accent"]])
                fig = apply_theme(fig)
                fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25))
                charts["mrr_trend"] = fig.to_html(full_html=False, include_plotlyjs="cdn")
        except Exception as e:
            print(f"MRR chart error: {e}")
    if rev_metrics.get("ltv_curve") and len(rev_metrics["ltv_curve"]) > 0:
        try:
            ltv_df = pd.DataFrame(rev_metrics["ltv_curve"])
            if len(ltv_df) > 0:
                ltv_df["percentile"] = pd.to_numeric(ltv_df["percentile"], errors="coerce").fillna(0).astype(int)
                ltv_df["ltv"] = pd.to_numeric(ltv_df["ltv"], errors="coerce").fillna(0)
                fig = px.line(ltv_df, x="percentile", y="ltv", template="plotly_dark")
                fig = apply_theme(fig)
                fig.update_traces(line_color=COLORS["blue"], fill="tozeroy", fillcolor="rgba(96,165,250,0.1)")
                fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25))
                charts["ltv_curve"] = fig.to_html(full_html=False, include_plotlyjs=False)
        except Exception as e:
            print(f"LTV chart error: {e}")
    if rev_metrics.get("subscription") and len(rev_metrics["subscription"]) > 0:
        try:
            sub_df = pd.DataFrame(rev_metrics["subscription"])
            if len(sub_df) > 0 and "revenue" in sub_df.columns and "subscription" in sub_df.columns:
                sub_df["revenue"] = pd.to_numeric(sub_df["revenue"], errors="coerce").fillna(0)
                fig = px.pie(sub_df, values="revenue", names="subscription", template="plotly_dark", hole=0.5)
                fig = apply_theme(fig)
                fig.update_traces(marker_colors=[COLORS["accent"], COLORS["blue"], COLORS["secondary"]], textinfo="label+percent", textfont=dict(size=11))
                fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=10), showlegend=False)
                charts["subscription"] = fig.to_html(full_html=False, include_plotlyjs=False)
        except Exception as e:
            print(f"Subscription chart error: {e}")
    if rev_metrics.get("cohort") and len(rev_metrics["cohort"]) > 0:
        try:
            cohort_df = pd.DataFrame(rev_metrics["cohort"])
            if len(cohort_df) > 0:
                pivot = cohort_df.pivot(index="cohort", columns="periods", values="retention").fillna(0)
                fig = px.imshow(pivot, template="plotly_dark", color_continuous_scale=["rgba(0,0,0,0)", COLORS["purple"]], text_auto=".0f")
                fig = apply_theme(fig)
                fig.update_layout(height=250, margin=dict(l=10, r=10, t=10, b=25), coloraxis_showscale=False)
                charts["cohort"] = fig.to_html(full_html=False, include_plotlyjs=False)
        except Exception as e:
            print(f"Cohort chart error: {e}")
    if rev_metrics.get("forecast") and len(rev_metrics["forecast"]) > 0:
        try:
            all_months = pd.DataFrame(rev_metrics["waterfall"] + rev_metrics["forecast"])
            all_months["revenue"] = pd.to_numeric(all_months["revenue"], errors="coerce").fillna(0)
            all_months["is_forecast"] = all_months.get("type", "").apply(lambda x: x == "forecast")
            fig = go.Figure()
            actual = all_months[~all_months["is_forecast"]]
            forecast = all_months[all_months["is_forecast"]]
            fig.add_trace(go.Scatter(x=actual["month"], y=actual["revenue"], mode="lines", name="Actual", line=dict(color=COLORS["primary"], width=2), fill="tozeroy", fillcolor="rgba(248,250,252,0.05)"))
            fig.add_trace(go.Scatter(x=forecast["month"], y=forecast["revenue"], mode="lines+markers", name="Forecast", line=dict(color=COLORS["accent"], width=2, dash="dash"), marker=dict(size=8)))
            fig = apply_theme(fig)
            fig.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=25), showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, font=dict(size=10)))
            charts["forecast"] = fig.to_html(full_html=False, include_plotlyjs=False)
        except Exception as e:
            print(f"Forecast chart error: {e}")
    return charts

def build_rfm_charts(rfm_data):
    charts = {}
    if not rfm_data:
        return charts
    rfm_df = pd.DataFrame(rfm_data)
    fig = px.scatter(rfm_df, x="recency", y="frequency", size="monetary", color="segment_name",
                   template="plotly_dark", color_discrete_sequence=COLORS["palette"])
    fig = apply_theme(fig)
    fig.update_traces(marker=dict(opacity=0.7, line=dict(width=0)))
    fig.update_layout(height=300, margin=dict(l=10, r=10, t=10, b=25), showlegend=True,
                      legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5, font=dict(size=10)))
    charts["rfm_scatter"] = fig.to_html(full_html=False, include_plotlyjs=False)

    seg_df = rfm_df["segment_name"].value_counts().reset_index()
    seg_df.columns = ["segment", "count"]
    fig2 = px.pie(seg_df, values="count", names="segment", template="plotly_dark", hole=0.5)
    fig2 = apply_theme(fig2)
    fig2.update_traces(marker_colors=COLORS["palette"], textinfo="label+percent", textfont=dict(size=11))
    fig2.update_layout(height=220, margin=dict(l=10, r=10, t=10, b=10), showlegend=False)
    charts["rfm_pie"] = fig2.to_html(full_html=False, include_plotlyjs=False)
    return charts

# ==================== SESSION STORE ====================
_session_data = {}

def get_session_data(session_id: str):
    return _session_data.get(session_id)

def set_session_data(session_id: str, df: pd.DataFrame, white_label: dict):
    metrics = calculate_metrics(df)
    rev_metrics = calculate_revenue_metrics(df)
    rfm = calculate_rfm(df)
    anomalies = detect_anomalies(df)
    _session_data[session_id] = {
        "df": df, "white_label": white_label,
        "metrics": metrics, "charts": build_charts(df, metrics),
        "rev_metrics": rev_metrics, "rev_charts": build_revenue_charts(df, rev_metrics),
        "rfm": rfm, "rfm_charts": build_rfm_charts(rfm),
        "anomalies": anomalies
    }

# ==================== ROUTES ====================
@app.get("/", response_class=HTMLResponse)
def index(request: Request):
    session_id = request.cookies.get("session_id", hashlib.sha256(str(datetime.now()).encode()).hexdigest()[:16])
    user = get_current_user(request)
    data = get_session_data(session_id)
    has_data = data is not None
    metrics = data["metrics"] if data else {}
    charts = data["charts"] if data else {}
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    response = templates.TemplateResponse("dashboard.html", {
        "request": request, "has_data": has_data, "metrics": metrics,
        "charts": charts, "white_label": white_label, "page": "dashboard", "user": user
    })
    response.set_cookie("session_id", session_id)
    return response

@app.get("/sources", response_class=HTMLResponse)
def sources_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("sources.html", {"request": request, "white_label": white_label, "page": "sources", "user": user})

@app.post("/load-sample")
def load_sample(request: Request):
    session_id = request.cookies.get("session_id")
    if not session_id:
        session_id = hashlib.sha256(str(datetime.now()).encode()).hexdigest()[:16]
    df = generate_sample_data()
    set_session_data(session_id, df, {"company_name": "MetricsAI", "logo_text": "MetricsAI"})
    response = RedirectResponse(url="/", status_code=303)
    response.set_cookie("session_id", session_id)
    return response

@app.post("/upload-csv")
async def upload_csv(request: Request, file: UploadFile = File(...)):
    session_id = request.cookies.get("session_id")
    if not session_id:
        session_id = hashlib.sha256(str(datetime.now()).encode()).hexdigest()[:16]
    contents = await file.read()
    try:
        if file.filename.endswith('.csv'):
            df = pd.read_csv(BytesIO(contents))
        elif file.filename.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(BytesIO(contents))
        else:
            return HTMLResponse("<div class='glass-card text-red-400 p-4'>❌ Поддерживаются только CSV и Excel</div>")
        for old, new in [('date', 'date'), ('user', 'user_id'), ('customer', 'user_id'), ('client', 'user_id'), ('revenue', 'revenue'), ('amount', 'revenue'), ('sum', 'revenue'), ('price', 'revenue'), ('value', 'revenue')]:
            col = next((c for c in df.columns if old in c.lower() and c != new), None)
            if col and col != new:
                df = df.rename(columns={col: new})
        if 'date' in df.columns:
            df['date'] = pd.to_datetime(df['date'])
        set_session_data(session_id, df, {"company_name": "MetricsAI", "logo_text": "MetricsAI"})
        conn = get_db()
        c = conn.cursor()
        user = get_current_user(request)
        user_id = user["id"] if user else None
        c.execute("INSERT INTO uploads (filename, source_type, upload_date, rows, cols, user_id) VALUES (?, ?, ?, ?, ?, ?)",
                  (file.filename, "upload", datetime.now().isoformat(), len(df), len(df.columns), user_id))
        conn.commit()
        conn.close()
        response = RedirectResponse(url="/", status_code=303)
        response.set_cookie("session_id", session_id)
        return response
    except Exception as e:
        return HTMLResponse(f"<div class='glass-card text-red-400 p-4'>❌ Ошибка: {str(e)}</div>")

@app.post("/connect-google-sheets")
async def connect_google_sheets(request: Request, url: str = Form(...)):
    session_id = request.cookies.get("session_id")
    if not session_id:
        session_id = hashlib.sha256(str(datetime.now()).encode()).hexdigest()[:16]
    try:
        csv_url = url.replace('/edit', '/export?format=csv') if '/edit' in url else (url + '/export?format=csv' if 'export?format=csv' not in url else url)
        df = pd.read_csv(csv_url)
        for old, new in [('date', 'date'), ('user', 'user_id'), ('revenue', 'revenue'), ('amount', 'revenue')]:
            col = next((c for c in df.columns if old in c.lower() and c != new), None)
            if col and col != new:
                df = df.rename(columns={col: new})
        if 'date' in df.columns:
            df['date'] = pd.to_datetime(df['date'])
        set_session_data(session_id, df, {"company_name": "MetricsAI", "logo_text": "MetricsAI"})
        conn = get_db()
        c = conn.cursor()
        user = get_current_user(request)
        user_id = user["id"] if user else None
        c.execute("INSERT INTO uploads (filename, source_type, upload_date, rows, cols, user_id) VALUES (?, ?, ?, ?, ?, ?)",
                  ("google_sheets", "google_sheets", datetime.now().isoformat(), len(df), len(df.columns), user_id))
        conn.commit()
        conn.close()
        response = RedirectResponse(url="/", status_code=303)
        response.set_cookie("session_id", session_id)
        return response
    except Exception as e:
        return HTMLResponse(f"<div class='glass-card text-red-400 p-4'>❌ Ошибка загрузки Google Sheets: {str(e)}</div>")

@app.post("/connect-postgresql")
async def connect_postgresql(request: Request, host: str = Form(...), port: str = Form("5432"), database: str = Form(...), user: str = Form(...), password: str = Form(...), query: str = Form(...)):
    session_id = request.cookies.get("session_id")
    if not session_id:
        session_id = hashlib.sha256(str(datetime.now()).encode()).hexdigest()[:16]
    try:
        import psycopg2
        conn = psycopg2.connect(host=host, port=port, database=database, user=user, password=password)
        df = pd.read_sql(query, conn)
        conn.close()
        if 'date' in df.columns:
            df['date'] = pd.to_datetime(df['date'])
        set_session_data(session_id, df, {"company_name": "MetricsAI", "logo_text": "MetricsAI"})
        response = RedirectResponse(url="/", status_code=303)
        response.set_cookie("session_id", session_id)
        return response
    except ImportError:
        return HTMLResponse("<div class='glass-card text-red-400 p-4'>❌ Установите psycopg2: <code>pip install psycopg2-binary</code></div>")
    except Exception as e:
        return HTMLResponse(f"<div class='glass-card text-red-400 p-4'>❌ Ошибка подключения: {str(e)}</div>")

@app.post("/api/v1/ingest")
async def api_ingest(request: Request):
    try:
        body = await request.json()
        session_id = request.headers.get("X-Session-ID", hashlib.sha256(str(datetime.now()).encode()).hexdigest()[:16])
        if isinstance(body, list):
            df = pd.DataFrame(body)
        elif isinstance(body, dict) and "data" in body:
            df = pd.DataFrame(body["data"])
        else:
            return JSONResponse({"error": "Expected JSON array or {data: [...]}"}, status_code=400)
        if 'date' in df.columns:
            df['date'] = pd.to_datetime(df['date'])
        set_session_data(session_id, df, {"company_name": "MetricsAI", "logo_text": "MetricsAI"})
        return JSONResponse({"success": True, "session_id": session_id, "rows": len(df), "cols": len(df.columns), "dashboard_url": f"/?session_id={session_id}"})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.get("/api/v1/health")
def api_health():
    return {"status": "ok", "version": "4.1", "timestamp": datetime.now().isoformat()}

@app.post("/share")
def share_dashboard(request: Request):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return JSONResponse({"error": "No data"}, status_code=400)
    dashboard_id = hashlib.sha256(str(datetime.now().timestamp()).encode()).hexdigest()[:12]
    conn = get_db()
    c = conn.cursor()
    c.execute("INSERT INTO shared_dashboards (id, created_at, data_json, white_label) VALUES (?, ?, ?, ?)",
              (dashboard_id, datetime.now().isoformat(), data["df"].to_json(), json.dumps(data["white_label"])))
    conn.commit()
    conn.close()
    return JSONResponse({"url": f"/dashboard/{dashboard_id}"})

@app.get("/dashboard/{dashboard_id}", response_class=HTMLResponse)
def public_dashboard(request: Request, dashboard_id: str):
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT data_json, white_label FROM shared_dashboards WHERE id = ?", (dashboard_id,))
    row = c.fetchone()
    conn.close()
    if not row:
        return templates.TemplateResponse("error.html", {"request": request, "message": "Dashboard not found"})
    df = pd.read_json(StringIO(row["data_json"]))
    white_label = json.loads(row["white_label"])
    metrics = calculate_metrics(df)
    charts = build_charts(df, metrics)
    return templates.TemplateResponse("public.html", {"request": request, "metrics": metrics, "charts": charts, "white_label": white_label})

@app.get("/funnel", response_class=HTMLResponse)
def funnel_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    has_data = data is not None
    charts = data["charts"] if data else {}
    metrics = data["metrics"] if data else {}
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("funnel.html", {"request": request, "has_data": has_data, "charts": charts, "metrics": metrics, "white_label": white_label, "page": "funnel", "user": user})

@app.get("/ai", response_class=HTMLResponse)
def ai_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    has_data = data is not None
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("ai.html", {"request": request, "has_data": has_data, "white_label": white_label, "page": "ai", "insights": "", "user": user})

@app.post("/ai-generate")
async def ai_generate(request: Request):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return HTMLResponse("<div class='glass-card text-red-400 p-4'>Сначала загрузите данные</div>")
    m = data["metrics"]
    df = data["df"]

    context_lines = [
        "Analytics Dashboard Data:",
        "- Total Revenue: {:,.0f} RUB".format(m.get('total_revenue', 0)),
        "- Total Users: {:,}".format(m.get('total_users', 0)),
        "- ARPU: {:,.0f} RUB".format(m.get('arpu', 0)),
        "- AOV: {:,.0f} RUB".format(m.get('aov', 0)),
        "- LTV: {:,.0f} RUB".format(m.get('ltv', 0)),
        "- LTV/CAC Ratio: {:.1f}x".format(m.get('ltv_cac_ratio', 0)),
        "- D30 Retention: {:.1f}%".format(m.get('d30_retention', 0)),
        "- Revenue Growth: {:+.1f}%".format(m.get('revenue_growth', 0)),
        "- User Growth: {:+.1f}%".format(m.get('user_growth', 0)),
        "- Best Channel: {}".format(m.get('best_channel', 'N/A')),
        "- Worst Channel: {}".format(m.get('worst_channel', 'N/A')),
        "- Overall Conversion: {:.2f}%".format(m.get('overall_conversion', 0)),
        "- Avg Session: {:.0f}s".format(m.get('avg_session_duration', 0)),
        "- Data Shape: {} rows, {} columns".format(df.shape[0], df.shape[1])
    ]
    context = "\n".join(context_lines)

    try:
        if OPENAI_API_KEY and OPENAI_API_KEY.startswith("sk-"):
            client = openai.OpenAI(api_key=OPENAI_API_KEY)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a senior product analyst. Analyze the data and provide 3-4 specific, actionable insights in Russian. Use HTML formatting with <b> tags for emphasis. Be concise and data-driven. Suggest specific next steps."},
                    {"role": "user", "content": context}
                ],
                temperature=0.7,
                max_tokens=800
            )
            ai_text = response.choices[0].message.content

            conn = get_db()
            c = conn.cursor()
            c.execute("INSERT INTO ai_conversations (session_id, user_message, ai_response) VALUES (?, ?, ?)",
                      (session_id, "generate_insights", ai_text))
            conn.commit()
            conn.close()

            html = '<div class="ai-insight"><div class="ai-header">🧠 AI Product Analyst (GPT-4)</div>'
            html += '<div style="color: #E2E8F0; line-height: 1.7; font-size: 0.95rem;">'
            html += ai_text.replace("\n", "<br>")
            html += '</div></div>'
            return HTMLResponse(html)
        else:
            raise Exception("No API key")
    except Exception as e:
        ltv_status = "ниже нормы" if m.get('ltv_cac_ratio', 0) < 3 else "в норме"
        retention_status = "низкая" if m.get('d30_retention', 0) < 20 else "хорошая"
        html = '<div class="ai-insight"><div class="ai-header">🧠 AI Product Analyst</div>'
        html += '<div style="color: #E2E8F0; line-height: 1.7;">'
        html += '<b>1. LTV/CAC ratio ({:.1f}x)</b> — {}.<br><br>'.format(m.get('ltv_cac_ratio', 0), ltv_status)
        html += '<b>2. ARPU {:,.0f} ₽</b> — анализируйте ARPU по каналам.<br><br>'.format(m.get('arpu', 0))
        html += '<b>3. D30 Retention {:.1f}%</b> — {}.<br><br>'.format(m.get('d30_retention', 0), retention_status)
        html += '<b>4. Лучший канал: {}</b>'.format(m.get('best_channel', 'N/A'))
        html += '</div></div>'
        return HTMLResponse(html)

@app.post("/ai-nlq")
async def ai_nlq(request: Request, query: str = Form(...)):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return HTMLResponse("<div class='glass-card text-red-400 p-4'>Сначала загрузите данные</div>")
    df = data["df"]
    columns = ", ".join(df.columns.tolist())

    try:
        if OPENAI_API_KEY and OPENAI_API_KEY.startswith("sk-"):
            client = openai.OpenAI(api_key=OPENAI_API_KEY)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "You are a data analyst. The user has a DataFrame with columns: {}. Based on their natural language query, determine the best visualization. Respond ONLY with a JSON object: {{'chart_type': 'line|bar|pie|scatter', 'x': 'column_name', 'y': 'column_name', 'groupby': 'column_name|null', 'title': 'Chart Title', 'aggregation': 'sum|mean|count'}}. No other text.".format(columns)},
                    {"role": "user", "content": query}
                ],
                temperature=0.3,
                max_tokens=200
            )
            import json as json_mod
            config = json_mod.loads(response.choices[0].message.content)
        else:
            raise Exception("No API key")

        chart_type = config.get("chart_type", "bar")
        x_col = config.get("x", df.columns[0])
        y_col = config.get("y", df.columns[1])
        groupby = config.get("groupby")
        title = config.get("title", "Chart")
        agg = config.get("aggregation", "sum")

        if groupby and groupby != "null" and groupby in df.columns:
            if agg == "sum":
                plot_df = df.groupby(groupby)[y_col].sum().reset_index()
            elif agg == "mean":
                plot_df = df.groupby(groupby)[y_col].mean().reset_index()
            else:
                plot_df = df.groupby(groupby)[y_col].count().reset_index()
        else:
            plot_df = df

        if chart_type == "line":
            fig = px.line(plot_df, x=x_col, y=y_col, template="plotly_dark")
            fig.update_traces(line_color=COLORS["primary"], line_width=2.5)
        elif chart_type == "pie":
            fig = px.pie(plot_df, values=y_col, names=x_col, template="plotly_dark", hole=0.5)
        elif chart_type == "scatter":
            fig = px.scatter(plot_df, x=x_col, y=y_col, template="plotly_dark")
            fig.update_traces(marker_color=COLORS["blue"])
        else:
            fig = px.bar(plot_df, x=x_col, y=y_col, template="plotly_dark")
            fig.update_traces(marker_color=COLORS["accent"], marker_opacity=0.8)

        fig = apply_theme(fig)
        fig.update_layout(height=350, margin=dict(l=10, r=10, t=40, b=25), title=dict(text=title, font_color=COLORS["primary"], font_size=14))
        chart_html = fig.to_html(full_html=False, include_plotlyjs="cdn")

        result_html = '<div class="glass-card animate-in"><div class="section-title">🤖 {}</div><div class="text-xs mb-3" style="color: var(--text-secondary);">Запрос: "{}"</div><div class="chart-container" style="height: 350px;">{}</div></div>'.format(title, query, chart_html)
        return HTMLResponse(result_html)
    except Exception as e:
        return HTMLResponse("<div class='glass-card text-red-400 p-4'>❌ Ошибка AI: {}</div>".format(str(e)))

@app.get("/ab", response_class=HTMLResponse)
def ab_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("ab.html", {"request": request, "white_label": white_label, "page": "ab", "user": user})

@app.post("/ab-calculate")
def ab_calculate(request: Request, conv_a: int = Form(120), n_a: int = Form(1000), conv_b: int = Form(150), n_b: int = Form(1000), confidence: float = Form(0.95)):
    p_a = conv_a / n_a if n_a > 0 else 0
    p_b = conv_b / n_b if n_b > 0 else 0
    se = np.sqrt(p_a * (1 - p_a) / n_a + p_b * (1 - p_b) / n_b)
    z = (p_b - p_a) / se if se > 0 else 0
    z_critical = 1.96 if confidence == 0.95 else 2.58
    is_significant = abs(z) > z_critical
    uplift = ((p_b - p_a) / p_a * 100) if p_a > 0 else 0
    html = '<div class="grid grid-cols-2 md:grid-cols-4 gap-4 mt-6"><div class="metric-container"><div class="metric-label">Конверсия A</div><div class="metric-value">{:.2%}</div></div><div class="metric-container"><div class="metric-label">Конверсия B</div><div class="metric-value">{:.2%}</div></div><div class="metric-container"><div class="metric-label">Uplift</div><div class="metric-value" style="color:{};">{:+.1f}%</div></div><div class="metric-container"><div class="metric-label">Значимость</div><div class="metric-value" style="color:{};">{}</div></div></div><div class="glass-card mt-4"><div style="color:#94A3B8;font-size:0.9rem">Z-score: <b>{:.3f}</b> | Критическое: ±{}</div></div>'.format(
        p_a, p_b, '#34D399' if uplift>0 else '#F87171', uplift, '#34D399' if is_significant else '#F87171', '✅ Да' if is_significant else '❌ Нет', z, z_critical)
    return HTMLResponse(html)

@app.get("/revenue", response_class=HTMLResponse)
def revenue_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    has_data = data is not None
    rev_metrics = data["rev_metrics"] if data else {}
    rev_charts = data["rev_charts"] if data else {}
    metrics = data["metrics"] if data else {}
    charts = data["charts"] if data else {}
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("revenue.html", {
        "request": request, "has_data": has_data, "rev_metrics": rev_metrics,
        "rev_charts": rev_charts, "metrics": metrics, "charts": charts, "white_label": white_label, "page": "revenue", "user": user
    })

@app.get("/builder", response_class=HTMLResponse)
def builder_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    has_data = data is not None
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("builder.html", {"request": request, "has_data": has_data, "white_label": white_label, "page": "builder", "user": user})

@app.post("/builder-generate")
def builder_generate(request: Request, metric: str = Form("revenue"), period: str = Form("monthly"), chart_type: str = Form("line"), layout: str = Form("standard")):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return HTMLResponse("<div class='glass-card text-red-400 p-4'>Сначала загрузите данные</div>")
    df = data["df"]
    df["date"] = pd.to_datetime(df["date"])
    if period == "monthly":
        df["period"] = df["date"].dt.strftime("%Y-%m")
    elif period == "weekly":
        df["period"] = df["date"].dt.strftime("%Y-W%W")
    else:
        df["period"] = df["date"].dt.strftime("%Y-%m-%d")
    grouped = df.groupby("period")[metric].sum().reset_index()

    if chart_type == "line":
        fig = px.line(grouped, x="period", y=metric, template="plotly_dark")
        fig.update_traces(line_color=COLORS["primary"], line_width=2.5, fill="tozeroy", fillcolor="rgba(248,250,252,0.05)")
    elif chart_type == "bar":
        fig = px.bar(grouped, x="period", y=metric, template="plotly_dark")
        fig.update_traces(marker_color=COLORS["secondary"], marker_opacity=0.6)
    elif chart_type == "area":
        fig = px.area(grouped, x="period", y=metric, template="plotly_dark")
        fig.update_traces(line_color=COLORS["accent"], fillcolor="rgba(52,211,153,0.1)")
    else:
        fig = px.scatter(grouped, x="period", y=metric, template="plotly_dark")
        fig.update_traces(marker_color=COLORS["blue"])

    fig = apply_theme(fig)
    fig.update_layout(height=400, margin=dict(l=10, r=10, t=40, b=25), title=dict(text="{} — {}".format(metric.upper(), period), font_color=COLORS["primary"], font_size=14))
    chart_html = fig.to_html(full_html=False, include_plotlyjs="cdn")
    return '<div class="glass-card"><h3 class="text-lg font-bold text-white mb-4">📊 {} ({}))</h3><div style="height: 400px; overflow: hidden;">{}</div></div>'.format(metric.upper(), period, chart_html)

@app.get("/rfm", response_class=HTMLResponse)
def rfm_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    has_data = data is not None
    rfm = data["rfm"] if data else []
    rfm_charts = data["rfm_charts"] if data else {}
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("rfm.html", {
        "request": request, "has_data": has_data, "rfm": rfm, "rfm_charts": rfm_charts,
        "white_label": white_label, "page": "rfm", "user": user
    })

@app.get("/alerts", response_class=HTMLResponse)
def alerts_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    has_data = data is not None
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT * FROM alerts WHERE user_id = ? ORDER BY created_at DESC", (user["id"] if user else -1,))
    alerts = [dict(row) for row in c.fetchall()]
    conn.close()
    return templates.TemplateResponse("alerts.html", {
        "request": request, "has_data": has_data, "alerts": alerts,
        "white_label": white_label, "page": "alerts", "user": user
    })

@app.post("/alerts/create")
async def create_alert(request: Request, metric: str = Form(...), threshold_type: str = Form(...), threshold_value: float = Form(...), webhook_url: str = Form("")):
    user = get_current_user(request)
    if not user:
        return JSONResponse({"error": "Unauthorized"}, status_code=401)
    conn = get_db()
    c = conn.cursor()
    c.execute("INSERT INTO alerts (user_id, metric, threshold_type, threshold_value, webhook_url, is_active) VALUES (?, ?, ?, ?, ?, 1)",
              (user["id"], metric, threshold_type, threshold_value, webhook_url))
    conn.commit()
    conn.close()
    return JSONResponse({"success": True, "message": "Alert created"})

@app.post("/alerts/check")
async def check_alerts(request: Request):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return JSONResponse({"error": "No data"}, status_code=400)
    metrics = data["metrics"]
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT * FROM alerts WHERE is_active = 1")
    alerts = c.fetchall()
    conn.close()
    triggered = []
    for alert in alerts:
        metric_val = metrics.get(alert["metric"], 0)
        threshold = alert["threshold_value"]
        if alert["threshold_type"] == "above" and metric_val > threshold:
            triggered.append({"alert": dict(alert), "value": metric_val})
        elif alert["threshold_type"] == "below" and metric_val < threshold:
            triggered.append({"alert": dict(alert), "value": metric_val})
    for t in triggered:
        if t["alert"]["webhook_url"]:
            try:
                import aiohttp
                async with aiohttp.ClientSession() as session:
                    payload = {
                        "text": "🚨 MetricsAI Alert: {} is {:.2f} (threshold: {} {})".format(t['alert']['metric'], t['value'], t['alert']['threshold_type'], t['alert']['threshold_value']),
                        "metric": t["alert"]["metric"],
                        "value": t["value"],
                        "timestamp": datetime.now().isoformat()
                    }
                    async with session.post(t["alert"]["webhook_url"], json=payload) as resp:
                        pass
            except:
                pass
    return JSONResponse({"triggered": len(triggered), "alerts": triggered})

# ==================== AUTH ROUTES ====================
@app.get("/login", response_class=HTMLResponse)
def login_page(request: Request):
    return templates.TemplateResponse("login.html", {"request": request, "error": None})

@app.post("/login")
async def login_post(request: Request, username: str = Form(...), password: str = Form(...)):
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT * FROM users WHERE username = ? OR email = ?", (username, username))
    user = c.fetchone()
    conn.close()
    if not user or not verify_password(password, user["hashed_password"]):
        return templates.TemplateResponse("login.html", {"request": request, "error": "Неверный логин или пароль"})
    access_token = create_access_token(data={"sub": user["username"]}, expires_delta=timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    response = RedirectResponse(url="/", status_code=303)
    response.set_cookie("access_token", access_token, httponly=True, max_age=ACCESS_TOKEN_EXPIRE_MINUTES * 60)
    return response

@app.get("/register", response_class=HTMLResponse)
def register_page(request: Request):
    return templates.TemplateResponse("register.html", {"request": request, "error": None})

@app.post("/register")
async def register_post(request: Request, username: str = Form(...), email: str = Form(...), password: str = Form(...), full_name: str = Form("")):
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT * FROM users WHERE username = ? OR email = ?", (username, email))
    if c.fetchone():
        conn.close()
        return templates.TemplateResponse("register.html", {"request": request, "error": "Пользователь уже существует"})
    hashed = get_password_hash(password)
    c.execute("INSERT INTO users (username, email, hashed_password, full_name) VALUES (?, ?, ?, ?)", (username, email, hashed, full_name))
    conn.commit()
    conn.close()
    access_token = create_access_token(data={"sub": username}, expires_delta=timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    response = RedirectResponse(url="/", status_code=303)
    response.set_cookie("access_token", access_token, httponly=True, max_age=ACCESS_TOKEN_EXPIRE_MINUTES * 60)
    return response

@app.post("/logout")
def logout():
    response = RedirectResponse(url="/login", status_code=303)
    response.delete_cookie("access_token")
    return response

# ==================== EXPORT ROUTES ====================
@app.get("/export", response_class=HTMLResponse)
def export_page(request: Request):
    user = get_current_user(request)
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    has_data = data is not None
    white_label = data["white_label"] if data else {"company_name": "MetricsAI", "logo_text": "MetricsAI"}
    return templates.TemplateResponse("export.html", {"request": request, "has_data": has_data, "white_label": white_label, "page": "export", "user": user})

@app.post("/export/png")
async def export_png(request: Request, chart_name: str = Form("revenue")):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return JSONResponse({"error": "No data"}, status_code=400)
    df = data["df"]
    metrics = data["metrics"]
    if chart_name == "revenue" and metrics.get("monthly_trend"):
        monthly_df = pd.DataFrame(metrics["monthly_trend"])
        fig = px.line(monthly_df, x="month", y="revenue", template="plotly_dark")
        fig = apply_theme(fig)
        fig.update_traces(line_color=COLORS["primary"], line_width=3)
    elif chart_name == "funnel" and metrics.get("funnel"):
        funnel_df = pd.DataFrame(metrics["funnel"])
        fig = px.funnel(funnel_df, x="users", y="step", template="plotly_dark")
        fig = apply_theme(fig)
    elif chart_name == "channel" and "channel" in df.columns:
        channel_df = df.groupby("channel")["revenue"].sum().reset_index().sort_values("revenue", ascending=True)
        fig = px.bar(channel_df, x="revenue", y="channel", orientation="h", template="plotly_dark")
        fig = apply_theme(fig)
    else:
        return JSONResponse({"error": "Chart not available"}, status_code=400)
    try:
        img_bytes = fig.to_image(format="png", width=1200, height=600, scale=2)
        return StreamingResponse(BytesIO(img_bytes), media_type="image/png", headers={"Content-Disposition": "attachment; filename={}.png".format(chart_name)})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/export/pdf")
async def export_pdf(request: Request):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return JSONResponse({"error": "No data"}, status_code=400)
    try:
        from weasyprint import HTML as WeasyHTML
        metrics = data["metrics"]
        charts = data["charts"]
        white_label = data["white_label"]
        html_content = '<!DOCTYPE html><html><head><style>body { font-family: "Inter", sans-serif; background: #000; color: #F5F5F7; padding: 40px; } .header { font-size: 28px; font-weight: 800; margin-bottom: 20px; } .grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 20px; margin-bottom: 30px; } .card { background: rgba(28,28,30,0.6); border: 1px solid rgba(255,255,255,0.08); border-radius: 16px; padding: 20px; } .label { font-size: 11px; color: #86868B; text-transform: uppercase; letter-spacing: 0.1em; } .value { font-size: 24px; font-weight: 700; margin-top: 8px; }</style></head><body><div class="header">{} — Dashboard Report</div><div class="grid"><div class="card"><div class="label">Revenue</div><div class="value">{:,.0f} ₽</div></div><div class="card"><div class="label">Users</div><div class="value">{:,}</div></div><div class="card"><div class="label">ARPU</div><div class="value">{:,.0f} ₽</div></div></div><div>{}</div></body></html>'.format(
            white_label['company_name'], metrics['total_revenue'], metrics['total_users'], metrics['arpu'], charts.get('revenue', ''))
        pdf = WeasyHTML(string=html_content).write_pdf()
        return StreamingResponse(BytesIO(pdf), media_type="application/pdf", headers={"Content-Disposition": "attachment; filename=dashboard.pdf"})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

@app.post("/export/pptx")
async def export_pptx(request: Request):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return JSONResponse({"error": "No data"}, status_code=400)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "{} — Analytics Report".format(data['white_label']['company_name'])
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(245, 245, 247)
        metrics = data["metrics"]
        y_pos = 1.8
        for label, value, color in [
            ("Revenue", "{:,.0f} ₽".format(metrics['total_revenue']), RGBColor(52, 211, 153)),
            ("Users", "{:,}".format(metrics['total_users']), RGBColor(96, 165, 250)),
            ("ARPU", "{:,.0f} ₽".format(metrics['arpu']), RGBColor(255, 159, 10)),
            ("LTV/CAC", "{:.1f}x".format(metrics['ltv_cac_ratio']), RGBColor(191, 90, 242)),
        ]:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(3), Inches(1.2))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(134, 134, 139)
            p2 = tf.add_paragraph()
            p2.text = value
            p2.font.size = Pt(28)
            p2.font.bold = True
            p2.font.color.rgb = color
            y_pos += 1.5
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(4), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Generated: {}".format(datetime.now().strftime('%Y-%m-%d %H:%M'))
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(99, 99, 102)
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return StreamingResponse(output, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": "attachment; filename=dashboard.pptx"})
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

# ==================== WEBSOCKET ====================
@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    await manager.connect(websocket)
    try:
        while True:
            data = await websocket.receive_text()
            msg = json.loads(data)
            if msg.get("action") == "ping":
                await websocket.send_json({"type": "pong", "timestamp": datetime.now().isoformat()})
            elif msg.get("action") == "subscribe_metrics":
                await websocket.send_json({"type": "subscribed", "channel": "metrics"})
    except WebSocketDisconnect:
        manager.disconnect(websocket)

@app.post("/simulate-live")
async def simulate_live(request: Request):
    session_id = request.cookies.get("session_id")
    data = get_session_data(session_id)
    if not data:
        return JSONResponse({"error": "No data"}, status_code=400)
    df = data["df"]
    new_row = {
        "date": datetime.now(),
        "user_id": "user_{:04d}".format(np.random.randint(0, 500)),
        "revenue": np.random.exponential(200) if np.random.random() > 0.7 else 0,
        "channel": np.random.choice(["organic", "instagram", "google", "facebook", "email", "tiktok", "referral"]),
        "event_type": np.random.choice(["impression", "click", "visit", "add_to_cart", "purchase"]),
        "marketing_cost": np.random.choice([0, 0, 0, 30, 50, 80, 120, 200]),
        "subscription": np.random.choice(["monthly", "annual", "none"], p=[0.15, 0.05, 0.80]),
        "session_duration": np.random.exponential(300),
        "page_views": np.random.poisson(5),
        "device": np.random.choice(["desktop", "mobile", "tablet"]),
        "country": np.random.choice(["RU", "US", "DE", "FR", "UK"]),
    }
    df = pd.concat([df, pd.DataFrame([new_row])], ignore_index=True)
    set_session_data(session_id, df, data["white_label"])
    await broadcast_metrics_update(session_id)
    return JSONResponse({"success": True, "new_rows": 1})

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
